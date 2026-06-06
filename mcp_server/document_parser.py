"""
文档解析层 — 统一解析 PPT/PDF/DOCX/图片/文本
支持格式：
  - PPT/PPTX   → python-pptx
  - PDF（电子） → PyMuPDF (fitz)
  - PDF（扫描） → pdfplumber + PaddleOCR
  - DOCX        → python-docx
  - 图片        → PaddleOCR
  - 纯文本      → 直接读取

输出统一结构：list[ParsedChunk]
"""

from pathlib import Path
from typing import List, Optional, Dict
from dataclasses import dataclass, field
import re
import os


@dataclass
class ParsedChunk:
    """解析后的文档块"""
    text: str                              # 文本内容
    source_file: str                       # 来源文件名
    source_type: str                       # ppt / pdf / docx / image / text
    page_number: Optional[int] = None      # 页码（PDF/PPT）
    slide_number: Optional[int] = None     # 幻灯片编号（PPT）
    chapter_hint: Optional[str] = None     # 章节提示（从标题自动提取）
    is_key_point: bool = False             # 是否为标记的重点
    extra_metadata: Dict = field(default_factory=dict)


# =========================================================================
# 主入口
# =========================================================================

def parse_file(file_path: str, is_key_point: bool = False) -> List[ParsedChunk]:
    """
    根据文件扩展名自动选择解析器，返回 ParsedChunk 列表。
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在：{file_path}")

    ext = path.suffix.lower()

    if ext in (".ppt", ".pptx"):
        chunks = _parse_ppt(file_path)
    elif ext == ".pdf":
        chunks = _parse_pdf(file_path)
    elif ext in (".doc", ".docx"):
        chunks = _parse_docx(file_path)
    elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif"):
        chunks = _parse_image(file_path)
    elif ext in (".txt", ".md", ".markdown", ".py", ".ipynb"):
        chunks = _parse_text(file_path)
    else:
        raise ValueError(f"不支持的文件格式：{ext}")

    # 标记重点
    if is_key_point:
        for c in chunks:
            c.is_key_point = True

    # 自动提取章节提示
    for c in chunks:
        c.chapter_hint = _guess_chapter(c.text)

    return chunks


def parse_directory(dir_path: str, recursive: bool = True) -> List[ParsedChunk]:
    """
    批量解析目录下所有支持的文件。
    """
    all_chunks = []
    supported_exts = {".ppt", ".pptx", ".pdf", ".doc", ".docx",
                      ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif",
                      ".txt", ".md", ".markdown"}

    dir_path = Path(dir_path)
    if not dir_path.exists():
        raise FileNotFoundError(f"目录不存在：{dir_path}")

    pattern = "**/*" if recursive else "*"
    for file_path in dir_path.glob(pattern):
        if file_path.suffix.lower() in supported_exts:
            try:
                chunks = parse_file(str(file_path))
                all_chunks.extend(chunks)
            except Exception as e:
                print(f"[WARN] 解析失败 {file_path}: {e}")

    return all_chunks


# =========================================================================
# PPT 解析器
# =========================================================================

def _parse_ppt(file_path: str) -> List[ParsedChunk]:
    """解析 PPT/PPTX 文件，每张幻灯片为一个 chunk"""
    chunks = []
    try:
        from pptx import Presentation
        prs = Presentation(file_path)

        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
                # 提取表格
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip():
                            texts.append(row_text)

            if texts:
                chunks.append(ParsedChunk(
                    text="\n".join(texts),
                    source_file=os.path.basename(file_path),
                    source_type="ppt",
                    slide_number=i,
                ))
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")

    return chunks


# =========================================================================
# PDF 解析器
# =========================================================================

def _parse_pdf(file_path: str) -> List[ParsedChunk]:
    """
    解析 PDF，优先使用 PyMuPDF（电子版），
    如果文本量过少则尝试 pdfplumber（扫描件处理提示）。
    """
    chunks = []
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)

        for i, page in enumerate(doc, start=1):
            text = page.get_text().strip()
            if text:
                chunks.append(ParsedChunk(
                    text=text,
                    source_file=os.path.basename(file_path),
                    source_type="pdf",
                    page_number=i,
                ))
        doc.close()

        # 如果 PyMuPDF 提取文本量很少，尝试 pdfplumber
        total_text = "".join(c.text for c in chunks)
        if len(total_text) < 100:
            _try_pdfplumber(file_path, chunks)

    except ImportError:
        # 回退到 pdfplumber
        _try_pdfplumber(file_path, chunks)

    return chunks


def _try_pdfplumber(file_path: str, chunks: List[ParsedChunk]) -> None:
    """使用 pdfplumber 解析（处理扫描件的文本层）"""
    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text()
                if text and text.strip():
                    # 避免重复添加
                    existing_texts = {c.text for c in chunks if c.page_number == i}
                    if text.strip() not in existing_texts:
                        chunks.append(ParsedChunk(
                            text=text.strip(),
                            source_file=os.path.basename(file_path),
                            source_type="pdf",
                            page_number=i,
                            extra_metadata={"parser": "pdfplumber"}
                        ))
    except ImportError:
        pass  # 静默跳过，后续可在 MCP 工具中提示用户安装


# =========================================================================
# DOCX 解析器
# =========================================================================

def _parse_docx(file_path: str) -> List[ParsedChunk]:
    """解析 Word 文档，按段落分组为 chunk"""
    chunks = []
    try:
        from docx import Document
        doc = Document(file_path)

        current_texts = []
        current_chapter = None

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                # 空行 = 自然分隔点
                if current_texts:
                    chunks.append(ParsedChunk(
                        text="\n".join(current_texts),
                        source_file=os.path.basename(file_path),
                        source_type="docx",
                        chapter_hint=current_chapter,
                    ))
                    current_texts = []
                continue

            # 检测标题样式（Word 中的 Heading 样式）
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                if current_texts:
                    chunks.append(ParsedChunk(
                        text="\n".join(current_texts),
                        source_file=os.path.basename(file_path),
                        source_type="docx",
                        chapter_hint=current_chapter,
                    ))
                    current_texts = []
                current_chapter = text
                # 标题单独作为一个 chunk
                chunks.append(ParsedChunk(
                    text=text,
                    source_file=os.path.basename(file_path),
                    source_type="docx",
                    chapter_hint=text,
                ))
                continue

            current_texts.append(text)

        # 收尾
        if current_texts:
            chunks.append(ParsedChunk(
                text="\n".join(current_texts),
                source_file=os.path.basename(file_path),
                source_type="docx",
                chapter_hint=current_chapter,
            ))

        # 解析表格
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_text.append(row_text)
            if table_text:
                chunks.append(ParsedChunk(
                    text="\n".join(table_text),
                    source_file=os.path.basename(file_path),
                    source_type="docx",
                    extra_metadata={"is_table": True}
                ))

    except ImportError:
        raise ImportError("请安装 python-docx: pip install python-docx")

    return chunks


# =========================================================================
# 图片 OCR 解析器
# =========================================================================

def _parse_image(file_path: str) -> List[ParsedChunk]:
    """使用 PaddleOCR 识别图片中的文字"""
    chunks = []
    try:
        from paddleocr import PaddleOCR
        import cv2
        import numpy as np

        # 初始化 OCR（首次运行下载模型）
        ocr = PaddleOCR(use_angle_cls=True, lang="ch", show_log=False)

        # 预处理：增强对比度
        img = cv2.imread(file_path)
        if img is None:
            raise ValueError(f"无法读取图片：{file_path}")

        # 转灰度 + 自适应二值化
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        enhanced = cv2.adaptiveThreshold(
            gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
            cv2.THRESH_BINARY, 11, 2
        )

        # 保存临时增强图片
        tmp_path = file_path + ".tmp_enhanced.png"
        cv2.imwrite(tmp_path, enhanced)

        result = ocr.ocr(tmp_path, cls=True)

        # 清理临时文件
        os.remove(tmp_path)

        if result and result[0]:
            texts = [line[1][0] for line in result[0] if line[1][0].strip()]
            if texts:
                chunks.append(ParsedChunk(
                    text="\n".join(texts),
                    source_file=os.path.basename(file_path),
                    source_type="image",
                    extra_metadata={"ocr_engine": "PaddleOCR"}
                ))

    except ImportError:
        # 回退方案：提示用户安装
        raise ImportError(
            "请安装 PaddleOCR: pip install paddlepaddle paddleocr\n"
            "或使用纯文本/PDF 格式替代图片。"
        )

    return chunks


# =========================================================================
# 纯文本解析器
# =========================================================================

def _parse_text(file_path: str) -> List[ParsedChunk]:
    """解析纯文本/Markdown 文件，按空行分隔为 chunk"""
    chunks = []
    path = Path(file_path)

    try:
        content = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        content = path.read_text(encoding="gbk", errors="replace")

    # 按空行分段
    paragraphs = re.split(r"\n\s*\n", content)

    for para in paragraphs:
        text = para.strip()
        if text:
            chunks.append(ParsedChunk(
                text=text,
                source_file=os.path.basename(file_path),
                source_type="text",
            ))

    return chunks


# =========================================================================
# 辅助函数
# =========================================================================

def _guess_chapter(text: str) -> Optional[str]:
    """
    从文本中自动识别章节标题。
    匹配模式如："第一章 xxx"、"第1章 xxx"、"Chapter 1" 等
    """
    patterns = [
        r"第[一二三四五六七八九十百千\d]+章\s*[^\n]*",
        r"Chapter\s*\d+\s*[^\n]*",
        r"第[一二三四五六七八九十百千\d]+节\s*[^\n]*",
        r"\d+[\.\、]\s*\S+",  # "1. xxx" 或 "1、xxx"
    ]

    lines = text.split("\n")[:3]  # 只看前3行
    for line in lines:
        line = line.strip()
        for pat in patterns:
            m = re.match(pat, line)
            if m:
                return m.group().strip()

    return None
